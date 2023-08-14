##Importaciones
import aspose.words as aw
from pptx import *
from pptx.util import Inches
import random

##Ruta 
ruta = "C:\\Users\colom\\OneDrive\\Documentos\\GitHub\\PDFHLocal\\HilosPDF\\Data\\Documents"

##Primo
def primo(num):
    for n in range(2, num):
        if num % n == 0:
            return False
    return True

##Documento Word
def documento_word(i):
    doc = aw.Document()
    builder = aw.DocumentBuilder(doc)
    contenido = "Hellor world" + str(i)
    builder.write(contenido)
    nombre = "out" + str(i) +".docx" 
    ruta_completa = ruta + "\\"+ nombre   
    print(ruta_completa)
    doc.save(ruta_completa)

##Crear PP
def powerpoint(i):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])  
    left = Inches(random.randint(1,5))
    top = Inches(random.randint(1,5))  
    width = Inches(random.randint(1,5))  
    height = Inches(random.randint(1,5))  
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = "hola mundo" +str(i)
    ruta_completa = ruta + "\\"+ f"out{i}.pptx"
    presentation.save(ruta_completa)

def crear_docs():
    for i in range(7):
        if(primo(i)):
            powerpoint(i)
        elif(i % 2 == 0):
            documento_word(i)

crear_docs()


